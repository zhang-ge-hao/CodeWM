import os
import sys
import json
import sys
import torch
import pickle
import random
from dataclasses import asdict
from typing import Tuple
from termcolor import colored
from tqdm import tqdm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
sys.path.append("src")

from _data import *
from _hf_obj import get_hf_tokenizer
from metrics import calculate
from _sweet import (
    WatermarkDetector
)

def tokenize(example: str, tokenizer) -> torch.Tensor:
    inputs = tokenizer(
        example,
        padding=True,
        truncation=True,
        return_tensors="pt",
        return_offsets_mapping=True
    )
    return inputs["input_ids"].squeeze(), inputs['offset_mapping'].squeeze()


def prompt_terminal(string):
    return colored(string, "grey")

def non_wm_token_terminal(string):
    return colored(string, "black")


def wm_token_terminal(string):
    return colored(string, color='green', attrs=['bold'])


def non_wm_token_ppt(text, paragraph):
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = 'Consolas'
    font.color.rgb = RGBColor(100, 100, 100)
    font.size = Pt(9)
    return run


def wm_token_ppt(text, paragraph):
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = 'Consolas'
    # font.color.rgb = RGBColor(0, 176, 80)
    font.color.rgb = RGBColor(0, 100, 0)
    font.bold = True
    font.underline = True
    font.size = Pt(9)
    return run


def detection(solution, model_name, gamma, custom_seed):
    tokenizer = get_hf_tokenizer(model_name)
    vocab = list(tokenizer.get_vocab().values())
    tokenized_suffix, offset_suffix = tokenize(solution, tokenizer)
    tokenized_text = tokenized_suffix
    detector = WatermarkDetector(vocab=vocab,
                                gamma=gamma,
                                tokenizer=tokenizer,
                                z_threshold=4.0,
                                ngram_len=5,
                                hash_key=custom_seed)
    detection_result_dict = detector.detect(tokenized_text=tokenized_text,
                                            tokenized_prefix=[])
    is_green_list = detection_result_dict["green_token_mask"]
    while len(is_green_list) < len(offset_suffix):
        is_green_list = [False] + is_green_list
    assert len(is_green_list) == len(offset_suffix), \
            f"{len(is_green_list)}, {len(offset_suffix)}"
    colored_case = ""
    for is_wm, (start, end) in zip(is_green_list, offset_suffix):
        token = solution[start: end]
        colored_case += wm_token_terminal(token) if is_wm else non_wm_token_terminal(token)
    detection_result_dict["colored_case"] = colored_case
    detection_result_dict["offset_suffix"] = offset_suffix
    return detection_result_dict, is_green_list, offset_suffix


with open("data/tmp/selected") as file:
    case_names = [l.strip() for l in file]

prs = Presentation()
for case_name in case_names:
    assert "wllm" in case_name
    obf_name, line_name, model_name, wm_name, para = case_name.split("--")
    ds_name = line_name.split("/")[0]

    model_2_shorter_name = {
        "meta-llama/Llama-3.1-8B-Instruct": "Llama31Instruct8B",
        "deepseek-ai/deepseek-coder-33b-base": "DSCoderBase33B",
    }
    result_path = f"data/result/{model_2_shorter_name[model_name]}--{wm_name}--{ds_name}"

    for fn2 in os.listdir(result_path):
        gp = f"{result_path}/{fn2}/generate.jsonl"
        op = f"{result_path}/{fn2}/obfuscate.jsonl"
        mp = f"{result_path}/{fn2}/metrics.jsonl"
        if "backup" in mp:
            continue
        with open(mp) as file:
            dps = [DataPoint(**json.loads(l)) for l in file]
        dp = dps[0]
        if para == f"{dp.temperature}/{dp.delta}/{dp.gamma}/{dp.entropy_threshold}":
            result_path = f"{result_path}/{fn2}"
            break

    with open(f"{result_path}/generate.jsonl") as file:
        gen_tasks = [GenTask(**json.loads(l)) for l in file]
        id_2_gen_task = {t.id: t for t in gen_tasks}
    with open(f"{result_path}/obfuscate.jsonl") as file:
        obf_tasks = [ObfTask(**json.loads(l)) for l in file]
        id_2_obf_task = {t.id: t for t in obf_tasks}

    obf_task = id_2_obf_task[case_name]
    gen_task = id_2_gen_task[obf_task.gen_task_id]

    for idx, task in enumerate([gen_task, obf_task]):
        detection_result, is_green_list, offset_suffix = detection(
            task.solution, 
            gen_task.model_name, gen_task.gamma, gen_task.custom_seed)
        print(detection_result["colored_case"])
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), 
                                        Inches(8), Inches(2))
        for is_wm, (start, end) in zip(is_green_list, offset_suffix):
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            token = task.solution[start: end]
            token = token.replace("\t", "    ")
            wm_token_ppt(token, p) if is_wm else non_wm_token_ppt(token, p)

prs.save(f"data/tmp/cases.pptx")
